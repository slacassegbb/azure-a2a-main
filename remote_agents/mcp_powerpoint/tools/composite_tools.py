"""
Composite tools for PowerPoint MCP Server.

These tools combine multiple operations into a single call to avoid
Azure AI Foundry's MCP chained tool call limitations.
"""
import os
import json
from typing import Dict, List, Optional, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from mcp.server.fastmcp import FastMCP
from mcp.types import ToolAnnotations


# ── Color schemes ────────────────────────────────────────────────────

COLOR_SCHEMES = {
    "modern_blue": {
        "primary": RGBColor(0x00, 0x78, 0xD4),
        "secondary": RGBColor(0x00, 0x5A, 0x9E),
        "accent": RGBColor(0x00, 0xB0, 0xF0),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light": RGBColor(0xE6, 0xF2, 0xFF),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "corporate_gray": {
        "primary": RGBColor(0x44, 0x44, 0x44),
        "secondary": RGBColor(0x66, 0x66, 0x66),
        "accent": RGBColor(0x00, 0x78, 0xD4),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light": RGBColor(0xF0, 0xF0, 0xF0),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "elegant_green": {
        "primary": RGBColor(0x2D, 0x7D, 0x46),
        "secondary": RGBColor(0x1B, 0x5E, 0x20),
        "accent": RGBColor(0x4C, 0xAF, 0x50),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light": RGBColor(0xE8, 0xF5, 0xE9),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "warm_red": {
        "primary": RGBColor(0xC6, 0x28, 0x28),
        "secondary": RGBColor(0xB7, 0x1C, 0x1C),
        "accent": RGBColor(0xFF, 0x8F, 0x00),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light": RGBColor(0xFF, 0xEB, 0xEE),
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
    },
}


def register_composite_tools(app: FastMCP, presentations: Dict, get_current_presentation_id):
    """Register composite tools with the FastMCP app."""

    @app.tool(
        annotations=ToolAnnotations(
            title="Build Presentation",
            destructiveHint=True,
        ),
    )
    def build_presentation(
        filename: str,
        slides: List[Dict[str, Any]],
        title: Optional[str] = None,
        author: Optional[str] = None,
        color_scheme: str = "modern_blue",
    ) -> Dict:
        """Create a complete PowerPoint presentation from a structured specification in a single call.

        This tool creates a new presentation and populates it with all slides in one
        operation.  Use it instead of calling create_presentation + add_slide +
        manage_text etc. individually.

        Args:
            filename: Name for the presentation file (with or without .pptx extension).
            slides: Ordered list of slide specs.  Each slide is a dict with a
                ``type`` key and type-specific fields.  Supported types:

                - ``{"type": "title", "title": "...", "subtitle": "..."}``
                - ``{"type": "content", "title": "...", "bullets": ["...", "..."]}``
                - ``{"type": "two_column", "title": "...", "left": ["..."], "right": ["..."]}``
                - ``{"type": "table", "title": "...", "headers": [...], "rows": [[...]]}``
                - ``{"type": "chart", "title": "...", "chart_type": "column", "categories": [...], "series": [{"name": "...", "values": [...]}]}``
                - ``{"type": "image", "title": "...", "image_source": "...", "source_type": "url|file"}``
                - ``{"type": "video", "title": "...", "video_source": "...", "source_type": "url|file"}``
                - ``{"type": "blank"}``

            title: Optional presentation metadata title.
            author: Optional presentation metadata author.
            color_scheme: Color theme - "modern_blue", "corporate_gray", "elegant_green", or "warm_red".

        Returns:
            JSON dict with status, filename, slide count, download_url, and any
            per-slide errors (the presentation is still saved even if individual
            slides fail).
        """
        if not filename.endswith(".pptx"):
            filename = f"{filename}.pptx"

        scheme = COLOR_SCHEMES.get(color_scheme, COLOR_SCHEMES["modern_blue"])

        try:
            pres = Presentation()

            # Set metadata
            if title:
                pres.core_properties.title = title
            if author:
                pres.core_properties.author = author

            results: List[Dict[str, Any]] = []

            for idx, slide_spec in enumerate(slides):
                slide_type = slide_spec.get("type", "").lower()
                try:
                    if slide_type == "title":
                        _add_title_slide(pres, slide_spec, scheme)
                    elif slide_type == "content":
                        _add_content_slide(pres, slide_spec, scheme)
                    elif slide_type == "two_column":
                        _add_two_column_slide(pres, slide_spec, scheme)
                    elif slide_type == "table":
                        _add_table_slide(pres, slide_spec, scheme)
                    elif slide_type == "chart":
                        _add_chart_slide(pres, slide_spec, scheme)
                    elif slide_type == "image":
                        _add_image_slide(pres, slide_spec, scheme)
                    elif slide_type in ("video", "video_slide"):
                        _add_video_slide(pres, slide_spec, scheme)
                    elif slide_type == "blank":
                        _add_blank_slide(pres)
                    else:
                        results.append({"index": idx, "type": slide_type, "ok": False, "error": f"Unknown slide type: {slide_type}"})
                        continue
                    results.append({"index": idx, "type": slide_type, "ok": True})
                except Exception as e:
                    results.append({"index": idx, "type": slide_type, "ok": False, "error": str(e)})

            # Save to download directory
            download_dir = "/tmp/pptx_downloads"
            os.makedirs(download_dir, exist_ok=True)
            safe_name = os.path.basename(filename).replace("/", "_").replace("\\", "_")
            if not safe_name.endswith(".pptx"):
                safe_name += ".pptx"
            dest_path = os.path.join(download_dir, safe_name)
            pres.save(dest_path)
            file_size = os.path.getsize(dest_path)

            # Store in presentations dict so other tools can access it
            pres_id = f"build_{safe_name}"
            presentations[pres_id] = pres

            errors = [r for r in results if not r.get("ok")]

            return {
                "message": "Presentation created and ready for download",
                "presentation_id": pres_id,
                "filename": safe_name,
                "slides_processed": len(results),
                "errors": errors if errors else None,
                "size_bytes": file_size,
                "download_url": f"/download/{safe_name}",
            }

        except Exception as e:
            return {"error": f"Failed to build presentation: {str(e)}"}


# ── Internal slide builders ──────────────────────────────────────────


def _get_blank_layout(pres):
    """Get the blank slide layout (index 6) or fall back to the last layout."""
    try:
        return pres.slide_layouts[6]
    except IndexError:
        return pres.slide_layouts[-1]


def _add_title_slide(pres, spec, scheme):
    """Add a title slide with large centered title and subtitle."""
    layout = _get_blank_layout(pres)
    slide = pres.slides.add_slide(layout)

    # Background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = scheme["primary"]

    # Title
    title_text = spec.get("title", "")
    left, top, width, height = Inches(0.5), Inches(2.0), Inches(9.0), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_text = spec.get("subtitle", "")
    if subtitle_text:
        left, top, width, height = Inches(1.0), Inches(3.8), Inches(8.0), Inches(1.0)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle_text
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER


def _add_content_slide(pres, spec, scheme):
    """Add a content slide with title and bullet points."""
    layout = _get_blank_layout(pres)
    slide = pres.slides.add_slide(layout)

    # Title bar
    _add_slide_title(slide, spec.get("title", ""), scheme)

    # Bullet points
    bullets = spec.get("bullets", [])
    if bullets:
        left, top, width, height = Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.font.color.rgb = scheme["text"]
            p.space_after = Pt(8)
            p.level = 0


def _add_two_column_slide(pres, spec, scheme):
    """Add a two-column layout slide."""
    layout = _get_blank_layout(pres)
    slide = pres.slides.add_slide(layout)

    _add_slide_title(slide, spec.get("title", ""), scheme)

    # Left column
    left_items = spec.get("left", [])
    if left_items:
        left, top, width, height = Inches(0.5), Inches(1.8), Inches(4.2), Inches(4.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(left_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = scheme["text"]
            p.space_after = Pt(6)

    # Right column
    right_items = spec.get("right", [])
    if right_items:
        left, top, width, height = Inches(5.3), Inches(1.8), Inches(4.2), Inches(4.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(right_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = scheme["text"]
            p.space_after = Pt(6)


def _add_table_slide(pres, spec, scheme):
    """Add a slide with a data table."""
    layout = _get_blank_layout(pres)
    slide = pres.slides.add_slide(layout)

    _add_slide_title(slide, spec.get("title", ""), scheme)

    headers = spec.get("headers", [])
    rows_data = spec.get("rows", [])
    total_rows = len(rows_data) + (1 if headers else 0)
    cols = len(headers) if headers else (len(rows_data[0]) if rows_data else 1)

    if total_rows == 0 or cols == 0:
        return

    left, top, width, height = Inches(0.5), Inches(1.8), Inches(9.0), Inches(0.4 * total_rows)
    table = slide.shapes.add_table(total_rows, cols, left, top, width, height).table

    # Header row
    if headers:
        for j, header in enumerate(headers):
            if j < cols:
                cell = table.cell(0, j)
                cell.text = str(header)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(11)
                    paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid()
                cell.fill.fore_color.rgb = scheme["primary"]

    # Data rows
    start_row = 1 if headers else 0
    for i, row_data in enumerate(rows_data):
        for j, cell_text in enumerate(row_data):
            if j < cols and (start_row + i) < total_rows:
                cell = table.cell(start_row + i, j)
                cell.text = str(cell_text)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.color.rgb = scheme["text"]


def _add_chart_slide(pres, spec, scheme):
    """Add a slide with a chart."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    layout = _get_blank_layout(pres)
    slide = pres.slides.add_slide(layout)

    _add_slide_title(slide, spec.get("title", ""), scheme)

    chart_type_str = spec.get("chart_type", "column").lower()
    chart_type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "stacked_column": XL_CHART_TYPE.COLUMN_STACKED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "stacked_bar": XL_CHART_TYPE.BAR_STACKED,
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
        "stacked_area": XL_CHART_TYPE.AREA_STACKED,
    }
    xl_chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)

    categories = spec.get("categories", [])
    series_list = spec.get("series", [])

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for s in series_list:
        chart_data.add_series(s.get("name", "Series"), s.get("values", []))

    left, top, width, height = Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5)
    slide.shapes.add_chart(xl_chart_type, left, top, width, height, chart_data)


def _add_image_slide(pres, spec, scheme):
    """Add a slide with an image."""
    import tempfile
    import urllib.request

    layout = _get_blank_layout(pres)
    slide = pres.slides.add_slide(layout)

    _add_slide_title(slide, spec.get("title", ""), scheme)

    image_source = spec.get("image_source", "")
    source_type = spec.get("source_type", "file")
    temp_path = None

    try:
        if source_type == "url" or image_source.startswith(("http://", "https://")):
            url_path = image_source.split("?")[0]
            ext = os.path.splitext(url_path)[1] or ".png"
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                temp_path = tmp.name
            urllib.request.urlretrieve(image_source, temp_path)
            abs_path = temp_path
        else:
            abs_path = os.path.abspath(image_source)
            if not os.path.exists(abs_path):
                return  # Skip if image not found

        img_width = spec.get("width")
        left, top = Inches(1.5), Inches(1.8)
        if img_width:
            slide.shapes.add_picture(abs_path, left, top, width=Inches(float(img_width)))
        else:
            slide.shapes.add_picture(abs_path, left, top, width=Inches(7.0))
    finally:
        if temp_path and os.path.exists(temp_path):
            os.unlink(temp_path)


def _add_video_slide(pres, spec, scheme):
    """Add a slide with an embedded video."""
    import tempfile
    import urllib.request

    _MIME_MAP = {
        ".mp4": "video/mp4",
        ".mov": "video/quicktime",
        ".avi": "video/x-msvideo",
        ".wmv": "video/x-ms-wmv",
        ".webm": "video/webm",
    }

    layout = _get_blank_layout(pres)
    slide = pres.slides.add_slide(layout)

    _add_slide_title(slide, spec.get("title", ""), scheme)

    video_source = spec.get("video_source", "")
    source_type = spec.get("source_type", "file")
    temp_video = None
    temp_poster = None

    try:
        # Resolve video path
        if source_type == "url" or video_source.startswith(("http://", "https://")):
            url_path = video_source.split("?")[0]
            ext = os.path.splitext(url_path)[1] or ".mp4"
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                temp_video = tmp.name
            urllib.request.urlretrieve(video_source, temp_video)
            video_path = temp_video
        else:
            video_path = os.path.abspath(video_source)
            ext = os.path.splitext(video_path)[1] or ".mp4"
            if not os.path.exists(video_path):
                return

        mime_type = _MIME_MAP.get(ext.lower(), "video/mp4")

        # Optional poster frame
        poster_path = None
        poster_source = spec.get("poster_source", "")
        if poster_source:
            poster_type = spec.get("poster_type", "file")
            if poster_type == "url" or poster_source.startswith(("http://", "https://")):
                p_ext = os.path.splitext(poster_source.split("?")[0])[1] or ".png"
                with tempfile.NamedTemporaryFile(delete=False, suffix=p_ext) as tmp:
                    temp_poster = tmp.name
                urllib.request.urlretrieve(poster_source, temp_poster)
                poster_path = temp_poster
            else:
                p = os.path.abspath(poster_source)
                if os.path.exists(p):
                    poster_path = p

        vid_width = Inches(float(spec.get("width", 7.0)))
        vid_height = Inches(float(spec.get("height", 4.0)))
        left, top = Inches(1.5), Inches(1.8)

        slide.shapes.add_movie(
            video_path,
            left, top, vid_width, vid_height,
            poster_frame_image=poster_path,
            mime_type=mime_type,
        )
    finally:
        if temp_video and os.path.exists(temp_video):
            os.unlink(temp_video)
        if temp_poster and os.path.exists(temp_poster):
            os.unlink(temp_poster)


def _add_blank_slide(pres):
    """Add a blank slide."""
    layout = _get_blank_layout(pres)
    pres.slides.add_slide(layout)


def _add_slide_title(slide, title_text, scheme):
    """Add a styled title bar to a slide."""
    if not title_text:
        return

    # Title background shape
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(1.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = scheme["primary"]
    shape.line.fill.background()

    # Title text
    left, top, width, height = Inches(0.5), Inches(0.2), Inches(9.0), Inches(0.9)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT
